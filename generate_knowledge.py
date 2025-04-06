import os
import glob
from pptx import Presentation
from docx import Document
import openpyxl
from PyPDF2 import PdfReader


def extractTextFromPPTX(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extractTextFromDOCX(file_path):
    document = Document(file_path)
    return "\n".join([paragraph.text for paragraph in document.paragraphs])

def extractTextFromXLSX(file_path):
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    text = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text.append(f"Sheet: {sheet_name}")
        for row in sheet.iter_rows():
            row_text = []
            for cell in row:
                if cell.value is not None:
                    row_text.append(str(cell.value))
            if row_text:
                text.append(" | ".join(row_text))
    return "\n".join(text)

def extractTextFromPDF(file_path):
    reader = PdfReader(file_path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text())
    return "\n".join(text)

def collectKnowledge():
    knowledge_dir = "./knowledge"
    all_texts = []

    for file_path in glob.glob(os.path.join(knowledge_dir, "*.*")):
        file_name = os.path.basename(file_path)
        file_ext = os.path.splitext(file_name)[1].lower()

        try:
            if file_ext == ".pptx":
                text = extractTextFromPPTX(file_path)
            elif file_ext == ".docx":
                text = extractTextFromDOCX(file_path)
            elif file_ext == ".xlsx":
                text = extractTextFromXLSX(file_path)
            elif file_ext == ".pdf":
                text = extractTextFromPDF(file_path)
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()

            formatted_text = f"\n\n==== FILE: {file_name} ====\n{text}"
            all_texts.append(formatted_text)
            print(f"Processed: {file_name}")
        except Exception as e:
            print(f"Error processing {file_name}: {str(e)}")

    return "\n\n".join(all_texts)

knowledge_text = collectKnowledge()

with open('./knowledge_text.txt', 'w', encoding='utf-8') as file:
    file.write(knowledge_text)

print("Knowledge base has been successfully generated and saved to knowledge_text.txt")
